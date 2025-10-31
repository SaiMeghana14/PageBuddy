# components_gemini.py
"""
Upgraded components_gemini.py (Option C)
- Robust Gemini (google.generativeai) + Vertex AI support
- TTS/STT with GCP clients if available, else safe fallbacks
- Extractive fallback summarizer
- Flashcard/todo parsers with resilient JSON parsing
- Helpers used by app.py: render_avatar, export_to_pptx, fetch_url_text, etc.
- Logging for debugging (Render / Streamlit Cloud friendly)
"""

import os
import json
import re
import time
import logging
import base64
import requests
from io import BytesIO

# HTML parser
from bs4 import BeautifulSoup

# NLP
import nltk
from nltk.tokenize import sent_tokenize
nltk.download("punkt", quiet=True)

# Vectorizer for extractive fallback
try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    import numpy as np
    SCIPY_AVAILABLE = True
except Exception:
    SCIPY_AVAILABLE = False

# PPTX
try:
    from pptx import Presentation
    from pptx.util import Pt
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# Logging
logger = logging.getLogger("components_gemini")
if not logger.handlers:
    logging.basicConfig(level=logging.INFO)

# --------- Gemini / Vertex detection ----------
GEN_CLIENT = None
genai = None
vertexai = None

# Prefer google.generativeai if available
try:
    import google.generativeai as genai_mod
    genai = genai_mod
    GEN_CLIENT = "genai"
    # configure API key if in env
    _genai_key = os.getenv("GENAI_API_KEY") or os.getenv("GOOGLE_API_KEY")
    if _genai_key:
        try:
            genai.configure(api_key=_genai_key)
        except Exception:
            # some environments permit ADC so ignore
            logger.debug("genai.configure() raised (ignored)")
except Exception:
    genai = None

# Vertex AI fallback
if GEN_CLIENT is None:
    try:
        import vertexai
        from vertexai import language as vlang
        vertexai = vertexai
        GEN_CLIENT = "vertexai"
    except Exception:
        vertexai = None
        GEN_CLIENT = None

logger.info("GEN_CLIENT: %s", GEN_CLIENT)

# --------- Google Cloud TTS/STT detection ----------
GCP_AUDIO = False
texttospeech = None
speech = None
try:
    from google.cloud import texttospeech as tts_mod, speech as speech_mod
    texttospeech = tts_mod
    speech = speech_mod
    GCP_AUDIO = True
except Exception:
    GCP_AUDIO = False

# --------- Helper: write service account json from secrets (Streamlit) ----------
def load_service_account_from_streamlit_secrets(st_secrets):
    """
    Accepts st.secrets (dict-like) expected to contain: {"google": {"credentials": <json or dict>}}
    Writes to /tmp and sets GOOGLE_APPLICATION_CREDENTIALS. Returns True on success.
    """
    try:
        if not st_secrets or "google" not in st_secrets:
            return False
        creds = st_secrets["google"].get("credentials") if isinstance(st_secrets["google"], dict) else st_secrets["google"]
        if not creds:
            return False
        if isinstance(creds, str):
            cred_dict = json.loads(creds)
        else:
            cred_dict = creds
        path = "/tmp/gcp_pagebuddy_creds.json"
        with open(path, "w") as f:
            json.dump(cred_dict, f)
        os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = path
        logger.info("Wrote GCP credentials to %s", path)
        # If genai present, try configure (ADC will be used by GCP libs)
        try:
            if genai:
                try:
                    genai.configure()  # no args will use ADC if available
                except Exception:
                    logger.debug("genai.configure() failed after writing ADC")
        except Exception:
            pass
        return True
    except Exception as e:
        logger.exception("load_service_account_from_streamlit_secrets failed: %s", e)
        return False

# --------- Fetch readable text from URL ----------
def fetch_url_text(url: str) -> str:
    """Return visible text from a URL (best-effort)."""
    try:
        headers = {"User-Agent": "PageBuddy/1.0 (+https://example.com)"}
        r = requests.get(url, timeout=8, headers=headers)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")
        for tag in soup(["script", "style", "noscript", "header", "footer", "form"]):
            tag.extract()
        text = soup.get_text(separator="\n", strip=True)
        # compress whitespace
        text = re.sub(r"\n{2,}", "\n\n", text)
        return text.strip()
    except Exception as e:
        logger.warning("fetch_url_text failed for %s: %s", url, e)
        return f"ERROR_FETCH: Could not fetch {url} ({e})"

# compatibility wrapper older code used
def extract_text_from_url(url: str) -> str:
    return fetch_url_text(url)

# --------- Gemini wrapper (unified) ----------
def _gemini_generate_text(prompt: str, model="gemini-1.5-flash", max_output_tokens=400, temperature=0.2, **kwargs):
    """
    Generate text using available client.
    Returns string or None on failure.
    """
    try:
        if GEN_CLIENT == "genai" and genai:
            try:
                # Use new high-level API: GenerativeModel
                model_obj = genai.GenerativeModel(model)
                response = model_obj.generate_content(
                    prompt,
                    **{"generation_config": {"max_output_tokens": int(max_output_tokens), "temperature": float(temperature)}}
                )
                # response.text is the common field
                text = getattr(response, "text", None)
                if text:
                    return str(text).strip()
                # fallback: try dict style
                if isinstance(response, dict):
                    return response.get("candidates", [{}])[0].get("content", "").strip()
                return str(response)
            except Exception as e:
                logger.exception("genai generate_content failed: %s", e)
                # Try older convenience API
                try:
                    resp = genai.generate_text(model=model, prompt=prompt, max_output_tokens=max_output_tokens, temperature=temperature)
                    if isinstance(resp, dict):
                        return resp.get("candidates",[{}])[0].get("content","").strip()
                    return getattr(resp, "text", getattr(resp, "content", str(resp))).strip()
                except Exception:
                    logger.exception("genai fallback generate_text failed")
                    return None
        elif GEN_CLIENT == "vertexai" and vertexai:
            try:
                from vertexai import language as vlang
                model_obj = vlang.TextGenerationModel.from_pretrained(model)
                response = model_obj.predict(prompt, max_output_tokens=max_output_tokens, temperature=temperature)
                return getattr(response, "text", None) or str(response)
            except Exception as e:
                logger.exception("vertexai generation failed: %s", e)
                return None
        else:
            logger.info("No Generative client available (GEN_CLIENT=%s)", GEN_CLIENT)
            return None
    except Exception as e:
        logger.exception("_gemini_generate_text unexpected error: %s", e)
        return None

# --------- Summarization & actions ----------
def smart_summarize(text: str, model="gemini-1.5-flash", language="English", style="anime"):
    """
    Primary summarization using Gemini; fallback to extractive summary.
    Returns text.
    """
    try:
        prompt = (
            f"You are NOVA, a calm futuristic assistant. Summarize the article into 4 short bullets, "
            f"then 3 concise action items and 5 short tags. Language: {language}. Style: {style}.\n\nArticle:\n{text[:16000]}"
        )
        out = _gemini_generate_text(prompt, model=model, max_output_tokens=420, temperature=0.12)
        if out and len(out.strip()) > 10:
            return out.strip()
    except Exception as e:
        logger.warning("smart_summarize primary failed: %s", e)
    # fallback
    return extractive_summary(text, n_sentences=6)

def generate_action_items(text: str, model="gemini-1.5-flash", language="English"):
    try:
        prompt = f"Create 4 concise action items from the text in {language}:\n\n{text[:12000]}"
        out = _gemini_generate_text(prompt, model=model, max_output_tokens=180, temperature=0.18)
        if out and len(out.strip())>5:
            return out.strip()
    except Exception as e:
        logger.debug("generate_action_items failed: %s", e)
    sents = sent_tokenize(text)
    return "\n".join(["- " + s.strip() for s in sents[:4]])

# --------- Extractive fallback ----------
def extractive_summary(text: str, n_sentences=6):
    """
    Lightweight extractive summarizer using TF-IDF sentence scoring.
    Falls back to truncation if sklearn not available.
    """
    try:
        sents = sent_tokenize(text)
        if len(sents) <= n_sentences:
            return "\n".join(sents)
        if SCIPY_AVAILABLE:
            vect = TfidfVectorizer(stop_words="english")
            X = vect.fit_transform(sents)
            scores = np.asarray(X.sum(axis=1)).ravel()
            ranked = np.argsort(scores)[::-1]
            top_idx = np.sort(ranked[:n_sentences])
            return " ".join([sents[i] for i in top_idx])
        else:
            # naive: pick first n sentences
            return " ".join(sents[:n_sentences])
    except Exception as e:
        logger.exception("extractive_summary error: %s", e)
        return text[:800] + ("..." if len(text) > 800 else "")

# --------- Flashcards / topics / todos ----------
def extract_topics(text: str, model="gemini-1.5-flash", top_n=6):
    try:
        prompt = f"Extract {top_n} short topics/section headers from this article, comma-separated:\n\n{text[:12000]}"
        out = _gemini_generate_text(prompt, model=model, max_output_tokens=160, temperature=0.0)
        if out:
            parts = re.split(r'[\n,;]+', out)
            return [p.strip() for p in parts if p.strip()][:top_n]
    except Exception:
        logger.debug("extract_topics failed")
    sents = sent_tokenize(text)
    return [s[:40] for s in sents[:top_n]]

def generate_flashcards(text: str, model="gemini-1.5-flash", count=8, language="English"):
    """
    Tries to get JSON from model; if not, parse Q/A pairs heuristically.
    Returns list of {"q":..,"a":..}
    """
    try:
        prompt = (
            f"Create {count} concise flashcards (Q -> A) from the article below. Provide them as a JSON list of objects "
            f'like [{{"q":"...", "a":"..."}}]. Language: {language}.\n\n{text[:15000]}'
        )
        out = _gemini_generate_text(prompt, model=model, max_output_tokens=420, temperature=0.2)
        if out:
            # try to extract a JSON array from the output
            try:
                # find first '[' and last ']' to parse JSON blob
                start = out.find("[")
                end = out.rfind("]") + 1
                if start != -1 and end != -1:
                    blob = out[start:end]
                    data = json.loads(blob)
                    if isinstance(data, list):
                        return data[:count]
            except Exception:
                # clean-line parse fallback
                cards = []
                lines = [l.strip() for l in out.splitlines() if l.strip()]
                q, a = None, None
                for l in lines:
                    if re.match(r"^\s*Q[:\-\)]", l, flags=re.I):
                        q = re.split(r"[:\-\)]\s*", l, maxsplit=1)[1].strip()
                    elif re.match(r"^\s*A[:\-\)]", l, flags=re.I):
                        a = re.split(r"[:\-\)]\s*", l, maxsplit=1)[1].strip()
                    else:
                        # sometimes "Q. ..." or numbered lists
                        if l.lower().startswith("q "):
                            q = l.split(" ",1)[1].strip()
                        elif l.lower().startswith("a "):
                            a = l.split(" ",1)[1].strip()
                    if q and a:
                        cards.append({"q": q, "a": a})
                        q, a = None, None
                if cards:
                    return cards[:count]
    except Exception as e:
        logger.debug("generate_flashcards exception: %s", e)

    # last-resort extractive generation
    sents = sent_tokenize(text)
    cards = []
    for i in range(count):
        q = sents[i*2] if i*2 < len(sents) else f"Concept {i+1}"
        a = sents[i*2+1] if i*2+1 < len(sents) else ""
        cards.append({"q": q, "a": a})
    return cards

def generate_todos(text: str, model="gemini-1.5-flash", language="English"):
    try:
        prompt = f"From the article, generate 6 actionable to-do items. Language: {language}\n\n{text[:12000]}"
        out = _gemini_generate_text(prompt, model=model, max_output_tokens=180, temperature=0.15)
        if out:
            return [l.strip("-• \n\t\r") for l in out.splitlines() if l.strip()][:6]
    except Exception:
        logger.debug("generate_todos failed")
    # fallback
    return ["Save article","Summarize key points","Make flashcards","Find references","Share with a peer","Schedule review"]

# --------- Sentiment -> Emotion mapping ----------
def sentiment_of_text(text: str, model="gemini-1.5-flash"):
    try:
        prompt = f"Provide one-word sentiment (positive/neutral/negative) for the text:\n\n{text[:5000]}"
        out = _gemini_generate_text(prompt, model=model, max_output_tokens=32, temperature=0.0)
        if out:
            first = out.strip().splitlines()[0].lower()
            if "positive" in first: return "positive"
            if "negative" in first: return "negative"
    except Exception:
        logger.debug("sentiment_of_text failed")
    return "neutral"

def analyze_emotion(text: str):
    s = sentiment_of_text(text)
    mapping = {"positive": "happy", "neutral": "listening", "negative": "concerned"}
    return mapping.get(s, "listening")

# --------- Translate helper ----------
def translate_text(text: str, target_language="Hindi", model="gemini-1.5-flash"):
    try:
        prompt = f"Translate to {target_language}:\n\n{text[:12000]}"
        out = _gemini_generate_text(prompt, model=model, max_output_tokens=400, temperature=0.0)
        if out:
            return out.strip()
    except Exception:
        logger.debug("translate_text failed")
    return text

# --------- TTS (GCP preferred) and fallback ----------
def tts_create_audio_bytes(text: str, language_code="en-IN", voice_name=None):
    """
    Return bytes of MP3 (binary) or None on failure.
    Supports:
      - Google Cloud Text-to-Speech (if GCP_AUDIO and credentials)
      - server-side pyttsx3 fallback (if available)
    voice_name: optional voice selector (e.g. "en-US-Neural2-C")
    """
    # prefer GCP
    try:
        if GCP_AUDIO and os.getenv("GOOGLE_APPLICATION_CREDENTIALS"):
            client = texttospeech.TextToSpeechClient()
            synthesis_input = texttospeech.SynthesisInput(text=text)
            # choose voice params if provided
            if voice_name:
                voice = texttospeech.VoiceSelectionParams(name=voice_name, language_code=language_code)
            else:
                voice = texttospeech.VoiceSelectionParams(language_code=language_code, ssml_gender=texttospeech.SsmlVoiceGender.FEMALE)
            audio_config = texttospeech.AudioConfig(audio_encoding=texttospeech.AudioEncoding.MP3)
            response = client.synthesize_speech(input=synthesis_input, voice=voice, audio_config=audio_config)
            return response.audio_content
    except Exception as e:
        logger.warning("GCP TTS failed: %s", e)

    # fallback to pyttsx3 (server-side)
    try:
        import pyttsx3
        tmp = "/tmp/pagebuddy_tts.mp3"
        engine = pyttsx3.init()
        # note: pyttsx3 may have different voice options per host; skip voice_name mapping
        engine.save_to_file(text, tmp)
        engine.runAndWait()
        with open(tmp, "rb") as f:
            return f.read()
    except Exception as e:
        logger.debug("pyttsx3 fallback failed: %s", e)
    return None

# --------- STT (GCP speech -> fallback to SpeechRecognition + pydub) ----------
def stt_from_uploaded_bytes(audio_bytes: bytes, language="en-IN"):
    """
    Accepts raw bytes of an audio file (any container).
    Returns transcribed string or "ERROR_STT:..." on failure.
    """
    # GCP speech
    try:
        if GCP_AUDIO and os.getenv("GOOGLE_APPLICATION_CREDENTIALS"):
            client = speech.SpeechClient()
            audio = speech.RecognitionAudio(content=audio_bytes)
            # best-effort config (let GCP auto-detect audio type)
            config = speech.RecognitionConfig(encoding=speech.RecognitionConfig.AudioEncoding.ENCODING_UNSPECIFIED,
                                              language_code=language, sample_rate_hertz=16000)
            response = client.recognize(config=config, audio=audio)
            return " ".join([r.alternatives[0].transcript for r in response.results])
    except Exception as e:
        logger.warning("GCP STT failed: %s", e)

    # local fallback via speech_recognition + pydub
    try:
        import speech_recognition as sr
        from pydub import AudioSegment
        tmp_in = "/tmp/pagebuddy_in_audio"
        with open(tmp_in, "wb") as f:
            f.write(audio_bytes)
        aud = AudioSegment.from_file(tmp_in)
        out_wav = "/tmp/pagebuddy_out.wav"
        aud.export(out_wav, format="wav")
        r = sr.Recognizer()
        with sr.AudioFile(out_wav) as source:
            audio = r.record(source)
        return r.recognize_google(audio)
    except Exception as e:
        logger.debug("fallback STT failed: %s", e)
        return f"ERROR_STT:{e}"

# --------- simple PPTX export ----------
def export_to_pptx(title: str, bullets, actions, filename="pagebuddy_export.pptx"):
    """
    bullets, actions: sequences of strings
    Returns BytesIO (seeked to 0) or None if python-pptx not available.
    """
    try:
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not available")
            return None
        prs = Presentation()
        # title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        # bullets slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key points"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for b in bullets:
            p = body.add_paragraph()
            p.text = str(b)
            p.level = 0
            p.font.size = Pt(18)
        # actions slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Action items"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for a in actions:
            p = body.add_paragraph()
            p.text = str(a)
            p.level = 0
            p.font.size = Pt(18)
        bio = BytesIO()
        prs.save(bio)
        bio.seek(0)
        return bio
    except Exception as e:
        logger.exception("export_to_pptx failed: %s", e)
        return None

# --------- tiny lipsync helper (estimate durations) ----------
def estimate_audio_duration_seconds(text: str):
    chars = max(1, len(text))
    # heuristic: 14 chars/sec typical, tuned slightly
    return max(1.0, chars / 14.0)

# --------- Render avatar helper (simple HTML) ----------
def render_avatar(state="listening", width=160):
    """
    Returns HTML snippet rendering the NOVA avatar with id 'nova_avatar'.
    Streamlit code expects this and will place it with unsafe_allow_html=True.
    """
    state = state or "listening"
    # pick png/svg
    img = f"avatar/nova_{state}.png"
    # fallback to idle if not present (client will 404 silently)
    html = f"""
    <div style="display:flex;align-items:center;gap:12px;">
      <img id="nova_avatar" src="{img}" class="holo-avatar" width="{width}" style="border-radius:12px;"/>
      <div style="display:flex;flex-direction:column;">
        <div style="font-weight:700;color:#e7fbff">NOVA</div>
        <div style="font-size:12px;color:#bfeefd">Your hologram assistant</div>
      </div>
    </div>
    """
    return html

# --------- Utilities: safe JSON parse helper ----------
def safe_json_loads(maybe_json: str):
    try:
        return json.loads(maybe_json)
    except Exception:
        # try to extract JSON substring
        try:
            s = maybe_json.strip()
            start = s.find("{")
            end = s.rfind("}") + 1
            if start != -1 and end != -1 and end > start:
                return json.loads(s[start:end])
        except Exception:
            pass
    return None

# --------- End of file ----------
logger.info("components_gemini loaded. GCP_AUDIO=%s, GEN_CLIENT=%s, PPTX=%s, SCIPY=%s", GCP_AUDIO, GEN_CLIENT, PPTX_AVAILABLE, SCIPY_AVAILABLE)
